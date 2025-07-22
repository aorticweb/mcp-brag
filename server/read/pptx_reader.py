from typing import Iterator

from pptx import Presentation

from server.constants import CHUNK_CHARACTER_LIMIT
from server.read.reader import Reader, SourceType, TextChunk
from server.read.text_reader import _split_text_chunk


class PptxReader(Reader):
    file_path: str
    chunk_size_max: int

    def __init__(self, file_path: str, chunk_size_max: int = CHUNK_CHARACTER_LIMIT.value):
        self.file_path = file_path
        self.chunk_size_max = chunk_size_max

    def source_type(self) -> SourceType:
        return SourceType.LOCAL_PPTX_FILE

    def read(self) -> str:
        """
        Read the PPTX file and return all text content as a single string.

        Returns:
            str: The extracted text content from all slides and shapes
        """
        prs = Presentation(self.file_path)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text

    def read_iter(self) -> Iterator[TextChunk]:
        """
        Read the PPTX file and yield TextChunk objects for each text shape.

        Each text shape becomes a TextChunk with start_index and end_index representing
        the character positions in the concatenated text from all slides.

        Yields:
            TextChunk: Chunk containing shape text and its character indices
        """
        prs = Presentation(self.file_path)
        char_index = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_text = shape.text

                    # Add newline to match the behavior of the read() method
                    shape_with_newline = shape_text + "\n"

                    # Only process non-empty text shapes (after stripping)
                    if shape_text.strip():
                        # Create TextChunk for this shape
                        shape_chunk = TextChunk(
                            start_index=char_index,
                            end_index=char_index + len(shape_with_newline),
                            text=shape_text.strip(),
                        )

                        # Split the shape chunk if it exceeds the size limit
                        for chunk in _split_text_chunk(self.chunk_size_max, shape_chunk):
                            yield chunk

                    # Always advance the character index by the full shape text length including newline
                    char_index += len(shape_with_newline)
